"""
Document processing utilities for text extraction and chunking.

This module provides tools to extract text from various document formats
and split them into manageable chunks for processing.
"""

import tiktoken
from typing import List, Dict, Optional


class DocumentProcessor:
    """
    Service for processing documents and generating text chunks.
    
    Supports extraction from PDF, DOCX, TXT, XLSX, PPTX and other formats.
    Provides intelligent text chunking with configurable overlap.
    """
    
    @staticmethod
    def count_tokens(text: str, model: str = "gpt-4") -> int:
        """
        Count tokens in text using tiktoken encoding.
        
        Args:
            text: Text to count tokens from
            model: Model name for encoding reference (default: "gpt-4")
            
        Returns:
            int: Number of tokens in the text
            
        Note:
            Falls back to character-based approximation (1 token ≈ 4 chars)
            if tiktoken encoding fails.
        """
        try:
            encoding = tiktoken.encoding_for_model(model)
            return len(encoding.encode(text))
        except Exception:
            # Fallback: approximation (1 token ≈ 4 characters)
            return len(text) // 4
    
    @staticmethod
    def extract_text_from_pdf(file_path: str) -> str:
        """
        Extract text content from PDF file.
        
        Args:
            file_path: Path to the PDF file
            
        Returns:
            str: Extracted text with page markers
            
        Raises:
            Exception: If PDF extraction fails
        """
        try:
            import PyPDF2
        except ImportError:
            raise ImportError(
                "PyPDF2 is required for PDF processing. "
                "Install with: pip install sonika-langchain-bot[documents]"
            )
        
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text += f"\n--- Página {page_num + 1} ---\n{page_text}\n"
        except Exception as e:
            raise Exception(f"Error extracting text from PDF: {str(e)}")
        return text.strip()
    
    @staticmethod
    def extract_text_from_docx(file_path: str) -> str:
        """
        Extract text content from DOCX file.
        
        Args:
            file_path: Path to the DOCX file
            
        Returns:
            str: Extracted text
            
        Raises:
            Exception: If DOCX extraction fails
        """
        try:
            import docx
        except ImportError:
            raise ImportError(
                "python-docx is required for DOCX processing. "
                "Install with: pip install sonika-langchain-bot[documents]"
            )
        
        try:
            doc = docx.Document(file_path)
            text = "\n".join([
                paragraph.text 
                for paragraph in doc.paragraphs 
                if paragraph.text.strip()
            ])
            return text.strip()
        except Exception as e:
            raise Exception(f"Error extracting text from DOCX: {str(e)}")
    
    @staticmethod
    def extract_text_from_txt(file_path: str) -> str:
        """
        Extract text content from plain text file.
        
        Args:
            file_path: Path to the text file
            
        Returns:
            str: File content
            
        Note:
            Attempts UTF-8 encoding first, falls back to latin-1
        """
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read().strip()
        except UnicodeDecodeError:
            # Fallback to latin-1 if UTF-8 fails
            with open(file_path, 'r', encoding='latin-1') as file:
                return file.read().strip()
    
    @staticmethod
    def extract_text_from_xlsx(file_path: str) -> str:
        """
        Extract text content from Excel file.
        
        Args:
            file_path: Path to the Excel file
            
        Returns:
            str: Extracted text with sheet and row markers
            
        Raises:
            Exception: If Excel extraction fails
        """
        try:
            from openpyxl import load_workbook
        except ImportError:
            raise ImportError(
                "openpyxl is required for Excel processing. "
                "Install with: pip install sonika-langchain-bot[documents]"
            )
        
        try:
            workbook = load_workbook(file_path, data_only=True)
            text = ""
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"\n--- Hoja: {sheet_name} ---\n"
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([
                        str(cell) if cell is not None else "" 
                        for cell in row
                    ])
                    if row_text.strip():
                        text += row_text + "\n"
            
            return text.strip()
        except Exception as e:
            raise Exception(f"Error extracting text from Excel: {str(e)}")
    
    @staticmethod
    def extract_text_from_pptx(file_path: str) -> str:
        """
        Extract text content from PowerPoint file.
        
        Args:
            file_path: Path to the PowerPoint file
            
        Returns:
            str: Extracted text with slide markers
            
        Raises:
            Exception: If PowerPoint extraction fails
        """
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError(
                "python-pptx is required for PowerPoint processing. "
                "Install with: pip install sonika-langchain-bot[documents]"
            )
        
        try:
            prs = Presentation(file_path)
            text = ""
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"\n--- Diapositiva {slide_num} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
            
            return text.strip()
        except Exception as e:
            raise Exception(f"Error extracting text from PowerPoint: {str(e)}")
    
    @classmethod
    def extract_text(cls, file_path: str, file_extension: str) -> str:
        """
        Extract text from file based on extension.
        
        Args:
            file_path: Path to the file
            file_extension: File extension (without dot)
            
        Returns:
            str: Extracted text
            
        Raises:
            ValueError: If file format is not supported
        """
        extractors = {
            'pdf': cls.extract_text_from_pdf,
            'docx': cls.extract_text_from_docx,
            'doc': cls.extract_text_from_docx,
            'txt': cls.extract_text_from_txt,
            'md': cls.extract_text_from_txt,
            'xlsx': cls.extract_text_from_xlsx,
            'xls': cls.extract_text_from_xlsx,
            'csv': cls.extract_text_from_txt,
            'pptx': cls.extract_text_from_pptx,
            'ppt': cls.extract_text_from_pptx,
        }
        
        extractor = extractors.get(file_extension.lower())
        if not extractor:
            supported = ', '.join(extractors.keys())
            raise ValueError(
                f"Format '{file_extension}' not supported. "
                f"Supported formats: {supported}"
            )
        
        return extractor(file_path)
    
    @classmethod
    def create_chunks(
        cls, 
        text: str, 
        chunk_size: int = 1000, 
        overlap: int = 200
    ) -> List[Dict]:
        """
        Split text into chunks with configurable overlap.
        
        Args:
            text: Complete text to chunk
            chunk_size: Maximum chunk size in tokens (default: 1000)
            overlap: Token overlap between chunks (default: 200)
            
        Returns:
            List[Dict]: List of chunks with metadata
                Each chunk contains:
                - content: Text content
                - chunk_index: Sequential index
                - token_count: Number of tokens
                - metadata: Additional metadata (empty dict)
                
        Example:
            >>> processor = DocumentProcessor()
            >>> chunks = processor.create_chunks("Long text...", chunk_size=500)
            >>> print(chunks[0])
            {
                'content': 'First chunk text...',
                'chunk_index': 0,
                'token_count': 450,
                'metadata': {}
            }
        """
        # Split into sentences (approximate)
        sentences = text.replace('\n', ' ').split('. ')
        
        chunks = []
        current_chunk = ""
        current_tokens = 0
        chunk_index = 0
        
        for sentence in sentences:
            sentence = sentence.strip()
            if not sentence:
                continue
            
            sentence_tokens = cls.count_tokens(sentence)
            
            # Check if adding sentence exceeds chunk size
            if current_tokens + sentence_tokens > chunk_size and current_chunk:
                # Save current chunk
                chunks.append({
                    'content': current_chunk.strip(),
                    'chunk_index': chunk_index,
                    'token_count': current_tokens,
                    'metadata': {}
                })
                
                # Prepare next chunk with overlap
                overlap_text = (
                    ' '.join(current_chunk.split()[-overlap:]) 
                    if overlap > 0 else ""
                )
                current_chunk = (
                    overlap_text + " " + sentence 
                    if overlap_text else sentence
                )
                current_tokens = cls.count_tokens(current_chunk)
                chunk_index += 1
            else:
                # Add sentence to current chunk
                current_chunk += (
                    " " + sentence if current_chunk else sentence
                )
                current_tokens += sentence_tokens
        
        # Add last chunk if exists
        if current_chunk.strip():
            chunks.append({
                'content': current_chunk.strip(),
                'chunk_index': chunk_index,
                'token_count': current_tokens,
                'metadata': {}
            })
        
        return chunks