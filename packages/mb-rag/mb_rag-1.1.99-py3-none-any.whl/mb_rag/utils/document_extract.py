"""
Document Extraction Module

This module provides functionality for extracting text and metadata from various document types
including CSV, PowerPoint (PPT/PPTX), and other document formats. It complements the PDF extraction
functionality in pdf_extract.py.

Example Usage:
    ```python
    # Initialize CSV extractor
    csv_extractor = CSVExtractor()
    
    # Extract data from a CSV file
    docs = csv_extractor.extract_csv("data.csv")
    
    # Initialize PowerPoint extractor
    ppt_extractor = PowerPointExtractor()
    
    # Extract content from a PowerPoint file
    docs = ppt_extractor.extract_ppt("presentation.pptx")
    ```

Features:
    - CSV file extraction with metadata
    - PowerPoint (PPT/PPTX) extraction
    - Batch processing for multiple files
"""

import os
import pandas as pd
import importlib.util
from typing import List, Dict, Optional, Union, Any
from langchain_core.documents import Document

class CSVExtractor:
    """
    Class for extracting data from CSV files.
    
    This class provides methods for extracting content from CSV files
    and converting it to Document objects for use with RAG systems.
    
    Args:
        logger: Optional logger instance for logging operations
        
    Example:
        ```python
        extractor = CSVExtractor()
        docs = extractor.extract_csv("data.csv")
        ```
    """
    
    def __init__(self, logger=None):
        """Initialize the CSV extractor."""
        self.logger = logger
    
    def check_file(self, file_path: str) -> bool:
        """
        Check if file exists.
        
        Args:
            file_path (str): Path to the file
            
        Returns:
            bool: True if file exists, False otherwise
        """
        return os.path.exists(file_path)
    
    def extract_csv(self, csv_path: str, include_stats: bool = True, 
                   chunk_by_row: bool = False, rows_per_chunk: int = 10,
                   **kwargs) -> List[Document]:
        """
        Extract data from a CSV file.
        
        Args:
            csv_path (str): Path to the CSV file
            include_stats (bool): Whether to include basic statistics in the metadata
            chunk_by_row (bool): Whether to create a separate document for each row or group of rows
            rows_per_chunk (int): Number of rows per chunk if chunk_by_row is True
            **kwargs: Additional arguments for pandas.read_csv
            
        Returns:
            List[Document]: List of Document objects containing extracted content
            
        Raises:
            ValueError: If the file doesn't exist
            ImportError: If pandas is not installed
        """
        if not self.check_file(csv_path):
            raise ValueError(f"File {csv_path} not found")
        
        try:
            # Read CSV file
            df = pd.read_csv(csv_path, **kwargs)
            
            # Create metadata
            metadata = {
                "source": csv_path,
                "rows": len(df),
                "columns": list(df.columns),
                "file_type": "csv"
            }
            
            # Add basic statistics if requested
            if include_stats:
                stats = {}
                for column in df.columns:
                    if pd.api.types.is_numeric_dtype(df[column]):
                        stats[column] = {
                            "min": float(df[column].min()),
                            "max": float(df[column].max()),
                            "mean": float(df[column].mean()),
                            "median": float(df[column].median())
                        }
                metadata["statistics"] = stats
            
            documents = []
            
            if chunk_by_row:
                # Create a separate document for each chunk of rows
                for i in range(0, len(df), rows_per_chunk):
                    chunk = df.iloc[i:i+rows_per_chunk]
                    chunk_text = chunk.to_string(index=False)
                    
                    chunk_metadata = metadata.copy()
                    chunk_metadata["chunk"] = {
                        "start_row": i,
                        "end_row": min(i + rows_per_chunk - 1, len(df) - 1),
                        "total_rows": len(chunk)
                    }
                    
                    documents.append(Document(
                        page_content=chunk_text,
                        metadata=chunk_metadata
                    ))
            else:
                # Create a single document with all data
                text = df.to_string(index=False)
                documents.append(Document(
                    page_content=text,
                    metadata=metadata
                ))
            
            if self.logger:
                self.logger.info(f"Extracted data from {csv_path}")
            else:
                print(f"Extracted data from {csv_path}")
                
            return documents
            
        except Exception as e:
            if self.logger:
                self.logger.error(f"Error extracting from {csv_path}: {str(e)}")
            else:
                print(f"Error extracting from {csv_path}: {str(e)}")
            raise
    
    def extract_multiple_csvs(self, csv_paths: List[str], **kwargs) -> List[Document]:
        """
        Extract data from multiple CSV files.
        
        Args:
            csv_paths (List[str]): List of paths to CSV files
            **kwargs: Additional arguments for extract_csv
            
        Returns:
            List[Document]: List of Document objects containing extracted content
        """
        all_docs = []
        for csv_path in csv_paths:
            try:
                docs = self.extract_csv(csv_path, **kwargs)
                all_docs.extend(docs)
                if self.logger:
                    self.logger.info(f"Successfully extracted content from {csv_path}")
                else:
                    print(f"Successfully extracted content from {csv_path}")
            except Exception as e:
                if self.logger:
                    self.logger.error(f"Error extracting from {csv_path}: {str(e)}")
                else:
                    print(f"Error extracting from {csv_path}: {str(e)}")
        
        return all_docs


class PowerPointExtractor:
    """
    Class for extracting content from PowerPoint (PPT/PPTX) files.
    
    This class provides methods for extracting text, notes, and metadata
    from PowerPoint presentations.
    
    Args:
        logger: Optional logger instance for logging operations
        
    Example:
        ```python
        extractor = PowerPointExtractor()
        docs = extractor.extract_ppt("presentation.pptx")
        ```
    """
    
    def __init__(self, logger=None):
        """Initialize the PowerPoint extractor."""
        self.logger = logger
    
    @staticmethod
    def check_package(package_name: str) -> bool:
        """
        Check if a Python package is installed.
        
        Args:
            package_name (str): Name of the package to check
            
        Returns:
            bool: True if package is installed, False otherwise
        """
        return importlib.util.find_spec(package_name) is not None
    
    def check_file(self, file_path: str) -> bool:
        """
        Check if file exists.
        
        Args:
            file_path (str): Path to the file
            
        Returns:
            bool: True if file exists, False otherwise
        """
        return os.path.exists(file_path)
    
    def extract_ppt(self, ppt_path: str, include_notes: bool = True,
                   include_hidden_slides: bool = False,
                   extract_images: bool = False) -> List[Document]:
        """
        Extract content from a PowerPoint file.
        
        Args:
            ppt_path (str): Path to the PowerPoint file
            include_notes (bool): Whether to include speaker notes
            include_hidden_slides (bool): Whether to include hidden slides
            extract_images (bool): Whether to extract images
            
        Returns:
            List[Document]: List of Document objects containing extracted content
            
        Raises:
            ValueError: If the file doesn't exist
            ImportError: If python-pptx is not installed
        """
        if not self.check_file(ppt_path):
            raise ValueError(f"File {ppt_path} not found")
        
        if not self.check_package("pptx"):
            raise ImportError("python-pptx package not found. Please install: pip install python-pptx")
        
        from pptx import Presentation
        
        try:
            # Load presentation
            presentation = Presentation(ppt_path)
            
            documents = []
            
            # Process each slide
            for i, slide in enumerate(presentation.slides):
                # Skip hidden slides if not requested
                if hasattr(slide, 'show') and not slide.show and not include_hidden_slides:
                    continue
                
                # Extract text from shapes
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        texts.append(shape.text)
                
                # Extract notes if requested
                notes = ""
                if include_notes and hasattr(slide, "notes_slide") and slide.notes_slide:
                    for note_shape in slide.notes_slide.notes_text_frame.paragraphs:
                        if note_shape.text:
                            notes += note_shape.text + "\n"
                
                # Create metadata
                metadata = {
                    "source": ppt_path,
                    "slide_number": i + 1,
                    "total_slides": len(presentation.slides),
                    "file_type": "pptx" if ppt_path.endswith(".pptx") else "ppt"
                }
                
                # Add slide title if available
                if slide.shapes.title and slide.shapes.title.text:
                    metadata["title"] = slide.shapes.title.text
                
                # Combine text content
                content = f"Slide {i+1}"
                if "title" in metadata:
                    content += f": {metadata['title']}"
                content += "\n\n"
                
                if texts:
                    content += "\n".join(texts) + "\n"
                
                if notes:
                    content += "\nNotes:\n" + notes
                
                # Create document
                documents.append(Document(
                    page_content=content,
                    metadata=metadata
                ))
            
            if self.logger:
                self.logger.info(f"Extracted {len(documents)} slides from {ppt_path}")
            else:
                print(f"Extracted {len(documents)} slides from {ppt_path}")
                
            return documents
            
        except Exception as e:
            if self.logger:
                self.logger.error(f"Error extracting from {ppt_path}: {str(e)}")
            else:
                print(f"Error extracting from {ppt_path}: {str(e)}")
            raise
    
    def extract_multiple_ppts(self, ppt_paths: List[str], **kwargs) -> List[Document]:
        """
        Extract content from multiple PowerPoint files.
        
        Args:
            ppt_paths (List[str]): List of paths to PowerPoint files
            **kwargs: Additional arguments for extract_ppt
            
        Returns:
            List[Document]: List of Document objects containing extracted content
        """
        all_docs = []
        for ppt_path in ppt_paths:
            try:
                docs = self.extract_ppt(ppt_path, **kwargs)
                all_docs.extend(docs)
                if self.logger:
                    self.logger.info(f"Successfully extracted content from {ppt_path}")
                else:
                    print(f"Successfully extracted content from {ppt_path}")
            except Exception as e:
                if self.logger:
                    self.logger.error(f"Error extracting from {ppt_path}: {str(e)}")
                else:
                    print(f"Error extracting from {ppt_path}: {str(e)}")
        
        return all_docs
