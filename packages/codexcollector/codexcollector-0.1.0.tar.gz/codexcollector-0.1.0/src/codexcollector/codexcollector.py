'''
Document ingestion system for extracting text and metadata from multiple formats.

This module provides the CodexCollector class, which serves as the entry point for document
collection in the LexicalForge pipeline. It handles both filesystem traversal and
web scraping, supporting multiple document formats including Word (.docx, .doc),
PowerPoint (.pptx, .ppt), PDF, and plain text files.

The module implements a Protocol-based architecture for text extraction, allowing
extensible support for additional formats through the TextExtractor protocol.

Key Features:
    - Automatic detection of file paths vs URLs
    - Recursive directory traversal with configurable exclusions
    - Web page scraping with document link discovery
    - Encoding detection for text files
    - Metadata extraction (creation dates) where available
    - Comprehensive error logging for failed ingestions
    - Rate limiting for web requests
    - File size validation

Typical Usage:
    collector = CodexCollector(max_file_size_mb=50, request_delay=1.5)
    codex = collector.collect('/path/to/documents')
    # or
    codex = collector.collect('https://example.com/docs')

The returned corpus (or codex) is a dictionary mapping integer IDs to document dictionaries
containing 'text', 'source', and 'date' keys.
'''

from pathlib import Path, PurePosixPath
from urllib.parse import urljoin, urlparse
from datetime import datetime
from typing import Protocol, TypedDict, Callable
import requests
import logging
import time
import chardet
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
import pypdf
import docx2txt
from .codex_exceptions import IngestionError, ConfigurationError

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


# Type definitions for better type safety
class Document(TypedDict):
    """Type definition for a document in the corpus."""
    filename: str
    source: str | None
    text: str
    date: str | None


# Type alias for corpus structure
CorpusType = dict[int, Document]


class TextExtractor(Protocol):
    '''Protocol for format-specific text extractors'''
    
    def extract_text(self, file_path: Path) -> str:
        '''Extract text content from file'''
        ...
    
    def extract_metadata(self, file_path: Path) -> dict[str, str | None]:
        '''Extract format-specific metadata'''
        ...

class WordExtractor:
    '''Handles .docx files using python-docx.
    
    Extracts text content and metadata (creation date) from Word documents.
    This extractor is designed for use within bulk operations, so it raises
    specific exceptions that allow callers to decide whether to fail fast
    or log and continue.
    
    Raises:
        IngestionError: When document cannot be opened or is corrupted.
        OSError: When file system access fails.
    '''
    
    def extract_text(self, file_path: Path) -> str:
        '''Extract text content from .docx file.
        
        Args:
            file_path: Path to the .docx file.
            
        Returns:
            Extracted text with paragraphs joined by newlines. Returns empty
            string if document contains no paragraphs.
            
        Raises:
            IngestionError: If document cannot be opened or is corrupted.
            OSError: If file cannot be accessed.
            
        Note:
            This method does not catch exceptions—callers must handle them
            according to their use case (fail fast vs. graceful degradation).
        '''
        try:
            doc = Document(file_path)
            return '\n'.join([paragraph.text for paragraph in doc.paragraphs])
        except (ValueError, KeyError) as e:
            # python-docx raises ValueError/KeyError for corrupted documents
            raise IngestionError(
                f"Failed to parse Word document {file_path}: corrupted or invalid format"
            ) from e
        except OSError as e:
            # File system errors (permission denied, file not found, etc.)
            raise OSError(
                f"Failed to access Word document {file_path}: {e}"
            ) from e
    
    def extract_metadata(self, file_path: Path) -> dict[str, str | None]:
        """Extract metadata from .docx file.
        
        Args:
            file_path: Path to the .docx file.
            
        Returns:
            Dictionary with 'date' key containing ISO format creation date
            or None if metadata is unavailable.
            
        Note:
            This method catches all exceptions and returns {'date': None} on
            failure, as metadata extraction is non-critical. The caller can
            still use the document even if metadata extraction fails.
        """
        try:
            doc = Document(file_path)
            core_props = doc.core_properties
            
            if core_props.created:
                return {'date': core_props.created.isoformat()}
            return {'date': None}
            
        except (ValueError, KeyError, OSError, AttributeError):
            # Metadata extraction is non-critical - return None on any failure
            # This allows text extraction to succeed even if metadata fails
            return {'date': None}

class LegacyWordExtractor:
    """Handles .doc files using docx2txt"""
    
    def extract_text(self, file_path: Path) -> str:
        return docx2txt.process(str(file_path))
    
    def extract_metadata(self, file_path: Path) -> dict[str, str | None]:
        # docx2txt doesn't provide metadata access for .doc files
        return {'date': None}

class PowerPointExtractor:
    """Handles .pptx and .ppt files using python-pptx.
    
    Extracts text content from all slides and shapes, along with metadata
    (creation date) from PowerPoint presentations. Designed for use within
    bulk operations with specific exception raising for caller control.
    
    Raises:
        IngestionError: When presentation cannot be opened or is corrupted.
        OSError: When file system access fails.
    """
    
    def extract_text(self, file_path: Path) -> str:
        """Extract text content from PowerPoint file.
        
        Iterates through all slides and shapes, collecting text from any
        shape that has a text attribute. Text from all shapes is joined
        with newlines.
        
        Args:
            file_path: Path to the .pptx or .ppt file.
            
        Returns:
            Extracted text with content from all slides joined by newlines.
            Returns empty string if presentation contains no text.
            
        Raises:
            IngestionError: If presentation cannot be opened or is corrupted.
            OSError: If file cannot be accessed.
            
        Note:
            python-pptx handles both .ppt and .pptx formats. This method
            does not catch exceptions—callers handle them based on context.
        """
        try:
            prs = Presentation(file_path)
            text_content = []
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content.append(shape.text)  # type: ignore
            
            return '\n'.join(text_content)
            
        except (ValueError, KeyError, AttributeError) as e:
            # python-pptx raises these for corrupted/invalid presentations
            raise IngestionError(
                f"Failed to parse PowerPoint {file_path}: corrupted or invalid format"
            ) from e
        except OSError as e:
            # File system errors (permission denied, file not found, etc.)
            raise OSError(
                f"Failed to access PowerPoint {file_path}: {e}"
            ) from e
    
    def extract_metadata(self, file_path: Path) -> dict[str, str | None]:
        """Extract metadata from PowerPoint file.
        
        Args:
            file_path: Path to the .pptx or .ppt file.
            
        Returns:
            Dictionary with 'date' key containing ISO format creation date
            or None if metadata is unavailable.
            
        Note:
            This method catches all exceptions and returns {'date': None} on
            failure, as metadata extraction is non-critical. Document can
            still be processed even if metadata extraction fails.
        """
        try:
            prs = Presentation(file_path)
            core_props = prs.core_properties
            
            if core_props.created:
                return {'date': core_props.created.isoformat()}
            return {'date': None}
            
        except (ValueError, KeyError, OSError, AttributeError):
            # Metadata extraction is non-critical - return None on any failure
            # Allows text extraction to succeed even if metadata fails
            return {'date': None}

class PDFExtractor:
    """Handles .pdf files using pypdf.
    
    Extracts text content from all pages and metadata (creation date) from
    PDF documents. Designed for use within bulk operations with specific
    exception raising for caller control.
    
    Raises:
        IngestionError: When PDF cannot be opened or is corrupted.
        OSError: When file system access fails.
    """
    
    def extract_text(self, file_path: Path) -> str:
        """Extract text content from PDF file.
        
        Iterates through all pages and extracts text using pypdf's
        text extraction capabilities. Text from all pages is joined
        with newlines.
        
        Args:
            file_path: Path to the .pdf file.
            
        Returns:
            Extracted text with content from all pages joined by newlines.
            Returns empty string if PDF contains no extractable text (e.g.,
            scanned documents without OCR).
            
        Raises:
            IngestionError: If PDF cannot be opened, is encrypted, or is corrupted.
            OSError: If file cannot be accessed.
            
        Note:
            pypdf does not perform OCR. Scanned PDFs without text layers
            will return empty strings. This method does not catch exceptions—
            callers handle them based on context.
        """
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = pypdf.PdfReader(file)
                text_content = []
                
                for page in pdf_reader.pages:
                    text_content.append(page.extract_text())
                
                return '\n'.join(text_content)
                
        except (ValueError, KeyError, RuntimeError) as e:
            # pypdf raises various exceptions for corrupted/encrypted PDFs
            # ValueError: malformed structures
            # KeyError: missing required PDF elements
            # RuntimeError: encrypted PDFs and other processing errors
            raise IngestionError(
                f"Failed to parse PDF {file_path}: corrupted, encrypted, or invalid format"
            ) from e
        except OSError as e:
            # File system errors (permission denied, file not found, etc.)
            raise OSError(
                f"Failed to access PDF {file_path}: {e}"
            ) from e
    
    def extract_metadata(self, file_path: Path) -> dict[str, str | None]:
        """Extract metadata from PDF file.
        
        Attempts to extract creation date from PDF metadata. PDF dates
        are in format D:YYYYMMDDHHmmSSOHH'mm and are parsed to ISO format.
        
        Args:
            file_path: Path to the .pdf file.
            
        Returns:
            Dictionary with 'date' key containing ISO format creation date
            or None if metadata is unavailable or cannot be parsed.
            
        Note:
            This method catches all exceptions and returns {'date': None} on
            failure, as metadata extraction is non-critical. Document can
            still be processed even if metadata extraction fails.
        """
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = pypdf.PdfReader(file)
                metadata = pdf_reader.metadata
                
                creation_date = None
                if metadata and '/CreationDate' in metadata:
                    # PDF metadata values are PdfObject - convert to string
                    date_obj = metadata['/CreationDate']
                    date_str = str(date_obj) if date_obj else None
                    
                    if date_str and date_str.startswith('D:'):
                        try:
                            # Extract YYYYMMDD portion
                            date_part = date_str[2:10]
                            creation_date = datetime.strptime(date_part, '%Y%m%d').isoformat()
                        except ValueError:
                            # Date string malformed - return None
                            pass
                
                return {'date': creation_date}
                
        except (ValueError, KeyError, RuntimeError, OSError, AttributeError):
            # Metadata extraction is non-critical - return None on any failure
            # Allows text extraction to succeed even if metadata fails
            return {'date': None}

class TextFileExtractor:
    """Handles plain text files (.txt, .md, .rtf).

    Extracts text content with automatic encoding detection using chardet.
    This extractor handles various text encodings to maximize compatibility
    with files from different sources and systems.

    Raises:
        IngestionError: When file cannot be decoded or is not valid text.
        OSError: When file system access fails.
    """

    def __init__(self, default_encoding: str = 'utf-8'):
        """Initialize TextFileExtractor with encoding configuration.

        Args:
            default_encoding: Encoding to use when chardet fails or as fallback.
        """
        self.default_encoding = default_encoding
    
    def extract_text(self, file_path: Path) -> str:
        """Extract text content from plain text file.
        
        Uses chardet to detect file encoding, then reads the file with the
        detected encoding. Falls back to UTF-8 if encoding detection fails.
        
        Args:
            file_path: Path to the text file.
            
        Returns:
            Extracted text content. Returns empty string if file is empty.
            
        Raises:
            IngestionError: If file cannot be decoded with detected encoding.
            OSError: If file cannot be accessed.
            
        Note:
            Encoding detection is not perfect and may fail on short files
            or files with mixed encodings. This method does not catch
            exceptions—callers handle them based on context.
        """
        try:
            # Detect encoding
            with open(file_path, 'rb') as file:
                raw_data = file.read()
                detected = chardet.detect(raw_data)
                encoding = detected['encoding'] or self.default_encoding
            
            # Read with detected encoding
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    return file.read()
            except (UnicodeDecodeError, LookupError) as e:
                # UnicodeDecodeError: encoding detection was wrong
                # LookupError: encoding name not recognized by Python
                raise IngestionError(
                    f"Failed to decode text file {file_path} with encoding {encoding}: {e}"
                ) from e
                
        except OSError as e:
            # File system errors (permission denied, file not found, etc.)
            raise OSError(
                f"Failed to access text file {file_path}: {e}"
            ) from e
    
    def extract_metadata(self, file_path: Path) -> dict[str, str | None]:
        """Extract metadata from text file.
        
        Plain text files do not contain embedded metadata, so this method
        always returns None for the date field.
        
        Args:
            file_path: Path to the text file.
            
        Returns:
            Dictionary with 'date' key set to None.
            
        Note:
            Unlike other extractors, this method never fails since it
            performs no actual operations. File modification times could
            be used but are handled at the CodexCollector level instead.
        """
        return {'date': None}

class CodexCollector:
    """
    Unified document collection system supporting filesystem and web sources.

    The CodexCollector class provides the primary interface for document ingestion in the
    LexicalForge pipeline. It automatically determines whether the source is a file
    path or URL and applies the appropriate collection strategy.

    For filesystem sources, CodexCollector recursively traverses directories, respecting
    exclusion patterns and file size limits. For web sources, it scrapes HTML pages
    to discover document links and retrieves them with configurable rate limiting.

    Attributes:
        supported_extensions: Set of file extensions to process during collection.
        max_file_size_mb: Maximum file size in MB (0 disables limit).
        request_delay: Delay in seconds between consecutive web requests.
        timeout: HTTP request timeout in seconds.
        excluded_dirs: Directory names to skip during filesystem traversal.
        error_log: Dictionary mapping source paths to error messages for failed ingestions.
        extractors: Dictionary mapping file extensions to TextExtractor implementations.
        corpus: Stored corpus from most recent collection (dict[int, dict[str, str]]).

    Methods:
        collect: Primary public method that ingests documents and returns corpus.

    Raises:
        ValueError: If source string is empty or invalid.
        FileNotFoundError: If filesystem path does not exist.
        requests.RequestException: If web requests fail after retries.

    Example:
        CodexCollector = CodexCollector(
            supported_extensions={'.pdf', '.docx', '.txt'},
            max_file_size_mb=50,
            excluded_dirs={'temp', 'cache'}
        )
        corpus = CodexCollector.collect('/data/documents')
        
        # Check for errors
        if CodexCollector.error_log:
            print(f"Failed to process {len(CodexCollector.error_log)} files")
        
        # Access stored corpus
        for doc_id, doc_data in CodexCollector.corpus.items():
            print(f"Doc {doc_id}: {len(doc_data['text'])} characters")
    """

    def __init__(
        self,
        supported_extensions: set[str] | None = None,
        max_file_size_mb: int = 100,
        request_delay: float = 1.0,
        timeout: int = 30,
        excluded_dirs: set[str] | None = None,
        max_crawl_depth: int = 2,
        default_encoding: str = 'utf-8',
        max_collection_time: int = 0
    ):
        """
        Initialize CodexCollector with configuration parameters.
        
        Args:
            supported_extensions: Set of file extensions to process. Must include
                leading dots (e.g., '.pdf', '.docx'). If None, uses default set.
            max_file_size_mb: Maximum file size in MB (0 = no limit). Must be
                non-negative.
            request_delay: Delay between web requests in seconds. Must be
                non-negative.
            timeout: HTTP request timeout in seconds. Must be positive.
            excluded_dirs: Directory names to skip during traversal. If None,
                uses default exclusion set.
            max_crawl_depth: Maximum depth for recursive web crawling (0 = only
                starting page, 1 = starting page + linked pages, etc.). Must be
                non-negative.
            default_encoding: Default text encoding to use when chardet fails
                or for performance. Common values: 'utf-8', 'latin-1', 'cp1252'.
            max_collection_time: Maximum time in seconds for entire collection
                (0 = no limit). Must be non-negative.
        
        Raises:
            ConfigurationError: If any parameter fails validation.
            
        Example:
            CodexCollector = CodexCollector(
                supported_extensions={'.pdf', '.docx', '.txt'},
                max_file_size_mb=50,
                excluded_dirs={'temp', 'cache'}
            )
        """
        # Validate max_file_size_mb
        if not isinstance(max_file_size_mb, int):
            raise ConfigurationError(
                f"max_file_size_mb must be int, got {type(max_file_size_mb).__name__}"
            )
        if max_file_size_mb < 0:
            raise ConfigurationError(
                f"max_file_size_mb must be non-negative, got {max_file_size_mb}"
            )
        
        # Validate request_delay
        if not isinstance(request_delay, (int, float)):
            raise ConfigurationError(
                f"request_delay must be numeric, got {type(request_delay).__name__}"
            )
        if request_delay < 0:
            raise ConfigurationError(
                f"request_delay must be non-negative, got {request_delay}"
            )
        
        # Validate timeout
        if not isinstance(timeout, int):
            raise ConfigurationError(
                f"timeout must be int, got {type(timeout).__name__}"
            )
        if timeout <= 0:
            raise ConfigurationError(
                f"timeout must be positive, got {timeout}"
            )

        # Validate max_crawl_depth
        if not isinstance(max_crawl_depth, int):
            raise ConfigurationError(
                f"max_crawl_depth must be int, got {type(max_crawl_depth).__name__}"
            )
        if max_crawl_depth < 0:
            raise ConfigurationError(
                f"max_crawl_depth must be non-negative, got {max_crawl_depth}"
            )

        # Validate default_encoding
        if not isinstance(default_encoding, str):
            raise ConfigurationError(
                f"default_encoding must be str, got {type(default_encoding).__name__}"
            )
        if not default_encoding or not default_encoding.strip():
            raise ConfigurationError(
                "default_encoding cannot be empty or whitespace-only"
            )

        # Validate max_collection_time
        if not isinstance(max_collection_time, int):
            raise ConfigurationError(
                f"max_collection_time must be int, got {type(max_collection_time).__name__}"
            )
        if max_collection_time < 0:
            raise ConfigurationError(
                f"max_collection_time must be non-negative, got {max_collection_time}"
            )
        
        # Validate and set supported_extensions
        if supported_extensions is not None:
            if not isinstance(supported_extensions, set):
                raise ConfigurationError(
                    f"supported_extensions must be set, got {type(supported_extensions).__name__}"
                )
            if not supported_extensions:
                raise ConfigurationError(
                    "supported_extensions cannot be empty set"
                )
            # Validate that all extensions start with dot
            invalid_extensions = {ext for ext in supported_extensions if not ext.startswith('.')}
            if invalid_extensions:
                raise ConfigurationError(
                    f"All extensions must start with '.', invalid: {invalid_extensions}"
                )
            self.supported_extensions = supported_extensions
        else:
            self.supported_extensions = {
                '.pdf', '.docx', '.doc', '.pptx', '.ppt', '.txt', '.md', '.rtf'
            }
        
        # Validate and set excluded_dirs
        if excluded_dirs is not None:
            if not isinstance(excluded_dirs, set):
                raise ConfigurationError(
                    f"excluded_dirs must be set, got {type(excluded_dirs).__name__}"
                )
            self.excluded_dirs = excluded_dirs
        else:
            self.excluded_dirs = {
                '.git', '__pycache__', '.pytest_cache', 'node_modules', '.vscode',
                '.idea', 'venv', '.env', 'logs', 'tmp', 'temp'
            }
        
        # Set validated attributes
        self.max_file_size_mb = max_file_size_mb
        self.request_delay = request_delay
        self.timeout = timeout
        self.max_crawl_depth = max_crawl_depth
        self.default_encoding = default_encoding
        self.max_collection_time = max_collection_time
        
        # Initialize error log for tracking failed ingestions
        self.error_log: dict[str, str] = {}
        
        # Initialize text extractors
        text_extractor = TextFileExtractor(default_encoding=self.default_encoding)
        self.extractors: dict[str, TextExtractor] = {
            '.docx': WordExtractor(),
            '.doc': LegacyWordExtractor(),
            '.pptx': PowerPointExtractor(),
            '.ppt': PowerPointExtractor(),  # python-pptx handles both .ppt and .pptx
            '.pdf': PDFExtractor(),
            '.txt': text_extractor,
            '.md': text_extractor,
            '.rtf': text_extractor,
        }

        # Corpus Attribute to store ingested documents
        self.corpus = None

    def _is_url(self, source: str) -> bool:
        """Determine if source string is a URL or file path.
        
        Uses simple heuristic based on protocol prefixes to distinguish between
        web URLs and filesystem paths. This is a lightweight check that avoids
        expensive parsing operations.
        
        Args:
            source: Source string to evaluate.
            
        Returns:
            True if source appears to be a URL, False if it appears to be a
            filesystem path.
            
        Raises:
            ValueError: If source is empty or not a string.
            
        Note:
            This method uses prefix matching rather than full URL parsing for
            performance. It recognizes http://, https://, and www. prefixes.
            Edge cases like 'ftp://' or relative URLs are treated as file paths.
            
        Example:
            self._is_url('https://example.com/docs')  # True
            self._is_url('www.example.com')           # True
            self._is_url('/home/user/documents')      # False
            self._is_url('documents/file.pdf')        # False
        """
        if not isinstance(source, str):
            raise ValueError(
                f"source must be string, got {type(source).__name__}"
            )
        
        if not source or not source.strip():
            raise ValueError("source cannot be empty or whitespace-only")
        
        return source.startswith(('http://', 'https://', 'www.'))
    
    def _is_valid_file(self, file_path: Path) -> bool:
        """Check if file meets collection criteria.
        
        Validates that a file should be included in collection based on extension
        and size criteria. This method is called during directory traversal to
        filter files before attempting extraction.
        
        Args:
            file_path: Path object representing the file to validate.
            
        Returns:
            True if file meets all criteria (supported extension, within size limit).
            False if file should be skipped.
            
        Note:
            This method logs validation failures to error_log for files that fail
            size checks or other validation issues (but not for unsupported extensions,
            which are silently skipped as expected behavior). Returns False rather
            than raising exceptions to allow directory traversal to continue.
            
        Example:
            if self._is_valid_file(Path('/docs/report.pdf')):
                # Process file
        """
        if not isinstance(file_path, Path):
            raise ValueError(
                f"file_path must be Path object, got {type(file_path).__name__}"
            )
        
        try:
            # Check extension
            file_ext = file_path.suffix.lower()
            if file_ext not in self.supported_extensions:
                return False
            
            # Check file size if limit is set
            if self.max_file_size_mb > 0:
                try:
                    file_size_mb = file_path.stat().st_size / (1024 * 1024)
                    if file_size_mb > self.max_file_size_mb:
                        logger.warning(
                            f"Skipping large file: {file_path} ({file_size_mb:.1f}MB "
                            f"exceeds limit of {self.max_file_size_mb}MB)"
                        )
                        return False
                except OSError as e:
                    # stat() can fail for permission issues, broken symlinks, etc.
                    self.error_log[str(file_path)] = f"File size check failed: {e}"
                    return False
            
            return True
            
        except OSError as e:
            # Unexpected OS errors (e.g., path too long, filesystem issues)
            self.error_log[str(file_path)] = f"File validation failed: {e}"
            return False
        except Exception as e:
            # Catch unexpected errors to prevent traversal failure
            # This ensures one problematic file doesn't stop entire collection
            self.error_log[str(file_path)] = f"Unexpected validation error: {e}"
            logger.error(f"Unexpected error validating {file_path}: {e}")
            return False
    
    def _should_skip_directory(self, dir_path: Path) -> bool:
        """Check if directory should be skipped during traversal"""
        return dir_path.name in self.excluded_dirs

    def _discover_files_from_path(self, path_source: str) -> list[str]:
        """Discover files from filesystem path, traversing subdirectories.
        
        Recursively traverses the directory tree starting from path_source,
        collecting all files that meet validation criteria (supported extension,
        size limits). Respects excluded_dirs to skip specified directories.
        
        Args:
            path_source: String path to file or directory. Can be absolute or relative.
            
        Returns:
            List of absolute file path strings for all discovered valid files.
            Returns empty list if path doesn't exist, no valid files found, or
            traversal fails.
            
        Note:
            This method logs errors to error_log but does not raise exceptions,
            allowing partial collection to succeed even if some paths fail.
            Permission errors on individual files/directories are logged and
            skipped rather than halting traversal.
            
        Example:
            files = CodexCollector._discover_files_from_path('/documents')
            # Returns: ['/documents/file1.pdf', '/documents/subdir/file2.docx']
        """
        if not isinstance(path_source, str):
            raise ValueError(
                f"path_source must be string, got {type(path_source).__name__}"
            )
        
        if not path_source or not path_source.strip():
            raise ValueError("path_source cannot be empty or whitespace-only")
        
        try:
            root = Path(path_source)
        except (ValueError, OSError) as e:
            # Path() can raise ValueError for invalid path strings (e.g., null bytes)
            # OSError for paths that are too long or contain invalid characters
            self.error_log[path_source] = f"Invalid path format: {e}"
            return []
        
        if not root.exists():
            self.error_log[path_source] = "Path does not exist"
            return []
        
        discovered_files = []
        
        try:
            if root.is_file():
                # Single file provided
                if self._is_valid_file(root):
                    discovered_files.append(str(root))
                return discovered_files
            
            # Directory traversal - recursive by requirement
            try:
                for file_path in root.rglob("*"):
                    try:
                        if file_path.is_file():
                            # Check if any parent directory should be skipped
                            skip = False
                            for parent in file_path.parents:
                                if parent == root:
                                    break
                                if self._should_skip_directory(parent):
                                    skip = True
                                    break
                            
                            if not skip and self._is_valid_file(file_path):
                                discovered_files.append(str(file_path))
                                
                    except PermissionError as e:
                        # Permission denied on specific file - log and continue
                        self.error_log[str(file_path)] = f"Permission denied: {e}"
                        continue
                    except OSError as e:
                        # Other OS errors on specific file (symlink loops, etc.)
                        self.error_log[str(file_path)] = f"File access error: {e}"
                        continue
                        
            except PermissionError as e:
                # Permission denied on directory traversal itself
                self.error_log[path_source] = f"Directory traversal permission denied: {e}"
                return discovered_files  # Return any files found before error
            except OSError as e:
                # Broader traversal errors (filesystem issues, too many open files, etc.)
                self.error_log[path_source] = f"Directory traversal failed: {e}"
                return discovered_files  # Return partial results
            
            return discovered_files
            
        except Exception as e:
            # Catch unexpected errors to prevent total collection failure
            self.error_log[path_source] = f"Unexpected discovery error: {e}"
            logger.error(f"Unexpected error during file discovery from {path_source}: {e}")
            return discovered_files  # Return any files discovered before failure

    def _discover_files_from_url(self, url_source: str) -> list[str]:
        """Discover document files from web URL.
        
        Crawls the starting URL and discovers linked document files within the
        same domain. Uses breadth-first traversal to find HTML pages and extracts
        links to supported document formats.
        
        Args:
            url_source: Starting URL for document discovery. Must include protocol
                (http:// or https://).
            
        Returns:
            List of document URLs discovered during crawling. Returns empty list
            if URL is inaccessible, crawling fails, or no documents are found.
            
        Note:
            This method logs errors to error_log but does not raise exceptions,
            allowing collection to continue even if web crawling fails. Respects
            request_delay between requests. Only crawls within the same domain
            as the starting URL.
            
        Example:
            urls = CodexCollector._discover_files_from_url('https://example.com/docs')
            # Returns: ['https://example.com/docs/file1.pdf', ...]
        """
        if not isinstance(url_source, str):
            raise ValueError(
                f"url_source must be string, got {type(url_source).__name__}"
            )
        
        if not url_source or not url_source.strip():
            raise ValueError("url_source cannot be empty or whitespace-only")
        
        # Validate URL format
        try:
            parsed = urlparse(url_source)
            if not parsed.scheme or not parsed.netloc:
                raise ValueError(
                    f"Invalid URL format: {url_source} (missing scheme or netloc)"
                )
            if parsed.scheme not in ('http', 'https'):
                raise ValueError(
                    f"Unsupported URL scheme: {parsed.scheme} (only http/https supported)"
                )
        except ValueError as e:
            self.error_log[url_source] = f"URL validation failed: {e}"
            return []
        
        discovered_urls = set()
        visited_urls = set()
        # Track URLs with their depth: (url, depth)
        urls_to_visit = [(url_source, 0)]
        
        try:
            base_domain = urlparse(url_source).netloc
        except Exception as e:
            self.error_log[url_source] = f"Failed to parse base domain: {e}"
            return []
        
        while urls_to_visit:
            current_url, current_depth = urls_to_visit.pop(0)

            if current_url in visited_urls:
                continue

            visited_urls.add(current_url)
            
            try:
                # Add delay between requests (skip for first request)
                if len(visited_urls) > 1:
                    time.sleep(self.request_delay)
                
                # Make HTTP request
                try:
                    response = requests.get(current_url, timeout=self.timeout)
                    response.raise_for_status()
                except requests.Timeout as e:
                    self.error_log[current_url] = f"Request timeout after {self.timeout}s: {e}"
                    continue
                except requests.HTTPError as e:
                    self.error_log[current_url] = f"HTTP error {response.status_code}: {e}"
                    continue
                except requests.ConnectionError as e:
                    self.error_log[current_url] = f"Connection failed: {e}"
                    continue
                except requests.RequestException as e:
                    self.error_log[current_url] = f"Request failed: {e}"
                    continue
                
                # Parse HTML content
                try:
                    soup = BeautifulSoup(response.content, 'html.parser')
                except Exception as e:
                    self.error_log[current_url] = f"HTML parsing failed: {e}"
                    continue
                
                # Find all links
                try:
                    for link in soup.find_all('a', href=True):
                        try:
                            href = link['href'].strip()
                            
                            if not href:
                                continue
                            
                            # Convert relative URLs to absolute
                            try:
                                full_url = urljoin(current_url, href)
                                parsed_url = urlparse(full_url)
                            except (ValueError, Exception) as e:
                                # urljoin or urlparse can fail on malformed URLs
                                logger.debug(f"Skipping malformed URL {href}: {e}")
                                continue
                            
                            # Skip non-HTTP protocols
                            if parsed_url.scheme not in ('http', 'https'):
                                continue
                            
                            # Stay within same domain
                            if parsed_url.netloc != base_domain:
                                continue
                            
                            # Check if it's a document file
                            url_path = parsed_url.path.lower()
                            if any(url_path.endswith(ext) for ext in self.supported_extensions):
                                discovered_urls.add(full_url)
                            # If within crawl depth, queue HTML pages for further crawling
                            elif current_depth < self.max_crawl_depth:
                                # Check if it looks like an HTML page
                                path_parts = parsed_url.path.split('/')
                                last_part = path_parts[-1] if path_parts else ''
                                # Queue if: ends with .html/.htm, is a directory, or has no extension
                                if (last_part.endswith(('.html', '.htm')) or
                                    not last_part or
                                    '.' not in last_part):
                                    if full_url not in visited_urls:
                                        urls_to_visit.append((full_url, current_depth + 1))
                                
                        except KeyError:
                            # link['href'] missing (shouldn't happen with href=True filter)
                            continue
                        except Exception as e:
                            # Catch unexpected errors processing individual links
                            logger.debug(f"Error processing link in {current_url}: {e}")
                            continue
                            
                except Exception as e:
                    self.error_log[current_url] = f"Link extraction failed: {e}"
                    continue
                    
            except Exception as e:
                # Catch-all for unexpected errors during URL processing
                self.error_log[current_url] = f"Web crawling failed: {e}"
                logger.error(f"Unexpected error crawling {current_url}: {e}")
                continue
        
        result_urls = list(discovered_urls)
        return result_urls

    def _check_collection_timeout(self, start_time: float) -> bool:
        """Check if collection has exceeded max_collection_time.

        Args:
            start_time: Timestamp when collection started.

        Returns:
            True if timeout exceeded (and max_collection_time > 0), False otherwise.
        """
        if self.max_collection_time <= 0:
            return False
        elapsed = time.time() - start_time
        return elapsed >= self.max_collection_time

    def _extract_filename_from_url(self, url: str) -> str:
        """Extract filename from URL, handling query strings and fragments.

        Uses URL parsing to extract the filename from the path component,
        properly handling query strings, fragments, and edge cases.

        Args:
            url: URL string to extract filename from.

        Returns:
            Filename extracted from URL path. Returns 'unknown' if URL is
            malformed or has no filename component.

        Example:
            self._extract_filename_from_url('https://example.com/doc.pdf?v=1')
            # Returns: 'doc.pdf'

            self._extract_filename_from_url('https://example.com/path/')
            # Returns: 'unknown'
        """
        try:
            parsed = urlparse(url)
            filename = PurePosixPath(parsed.path).name
            return filename if filename else 'unknown'
        except (ValueError, Exception):
            return 'unknown'

    def _download_file(self, url: str, temp_path: Path) -> None:
        """Download file from URL to temporary location.
        
        Streams file content from URL to local temporary file using chunked
        transfer to handle large files efficiently. This is a helper method
        for URL-based document collection.
        
        Args:
            url: Source URL to download from. Must be valid HTTP/HTTPS URL.
            temp_path: Path object indicating where to save downloaded file.
                Parent directory must exist and be writable.
        
        Raises:
            IngestionError: If download fails due to HTTP errors or invalid response.
            OSError: If file cannot be written to temp_path.
            ValueError: If url or temp_path are invalid.
            
        Note:
            This method raises exceptions rather than logging them, allowing
            the caller (collect method) to decide how to handle download failures
            in the context of bulk operations. Uses streaming to avoid loading
            entire file into memory.
            
        Example:
            temp_file = Path('temp_document.pdf')
            CodexCollector._download_file('https://example.com/doc.pdf', temp_file)
        """
        # Input validation
        if not isinstance(url, str):
            raise ValueError(
                f"url must be string, got {type(url).__name__}"
            )
        
        if not url or not url.strip():
            raise ValueError("url cannot be empty or whitespace-only")
        
        if not isinstance(temp_path, Path):
            raise ValueError(
                f"temp_path must be Path object, got {type(temp_path).__name__}"
            )
        
        # Validate parent directory exists
        try:
            if not temp_path.parent.exists():
                raise ValueError(
                    f"Parent directory does not exist: {temp_path.parent}"
                )
        except OSError as e:
            raise OSError(
                f"Cannot access parent directory {temp_path.parent}: {e}"
            ) from e
        
        try:
            # Make HTTP request with streaming enabled
            try:
                response = requests.get(url, stream=True, timeout=self.timeout)
                response.raise_for_status()

                # Log response information for debugging
                content_type = response.headers.get('content-type', 'unknown')
                content_length = response.headers.get('content-length', 'unknown')
                logger.debug(f"Downloading {url}: {content_type}, {content_length} bytes")

            except requests.Timeout as e:
                raise IngestionError(
                    f"Download timeout after {self.timeout}s for {url}"
                ) from e
            except requests.HTTPError as e:
                raise IngestionError(
                    f"HTTP error {response.status_code} downloading {url}: {e}"
                ) from e
            except requests.ConnectionError as e:
                raise IngestionError(
                    f"Connection failed for {url}: {e}"
                ) from e
            except requests.RequestException as e:
                raise IngestionError(
                    f"Request failed for {url}: {e}"
                ) from e
            
            # Write response content to file in chunks
            try:
                with open(temp_path, 'wb') as file:
                    for chunk in response.iter_content(chunk_size=8192):
                        if chunk:  # Filter out keep-alive chunks
                            file.write(chunk)
            except OSError as e:
                # File write failures (disk full, permissions, etc.)
                raise OSError(
                    f"Failed to write downloaded content to {temp_path}: {e}"
                ) from e
            except Exception as e:
                # Unexpected errors during streaming (connection dropped mid-transfer)
                raise IngestionError(
                    f"Download interrupted for {url}: {e}"
                ) from e
            
            # Validate downloaded file
            if not temp_path.exists():
                raise IngestionError(
                    f"Download completed but file not found: {temp_path}"
                )
            
            if temp_path.stat().st_size == 0:
                raise IngestionError(
                    f"Downloaded file is empty: {url}"
                )
                
        except (IngestionError, OSError, ValueError):
            # Re-raise expected exceptions without modification
            raise
        except Exception as e:
            # Catch unexpected errors and wrap in IngestionError
            raise IngestionError(
                f"Unexpected download error for {url}: {e}"
            ) from e

    def _extract_from_file(self, file_path: Path, source: str | None = None) -> Document:
        """Extract text and metadata from a single file.
        
        Uses the appropriate format-specific extractor to retrieve text content
        and metadata from the file. This method is designed for bulk operations
        and logs errors rather than raising them, allowing collection to continue
        even when individual files fail.
        
        Args:
            file_path: Path object pointing to the file to extract.
            source: Optional source identifier (URL or path string) for error logging.
                If None, uses file_path string representation.
            
        Returns:
            Dictionary with keys:
                - 'filename': Name of the file
                - 'source': Source identifier (URL or file path)
                - 'text': Extracted text content (empty string on failure)
                
        Note:
            This method never raises exceptions - all failures are logged to
            error_log and result in returning a document dict with empty text.
            This allows bulk collection to continue despite individual file failures.
            
        Example:
            result = CodexCollector._extract_from_file(Path('/docs/report.pdf'), '/docs/report.pdf')
            # Returns: {'filename': 'report.pdf', 'source': '/docs/report.pdf', 'text': '...'}
        """
        if not isinstance(file_path, Path):
            error_msg = f"file_path must be Path object, got {type(file_path).__name__}"
            file_key = source or str(file_path)
            self.error_log[file_key] = error_msg
            return {
                'filename': str(file_path) if source and self._is_url(source) else 'unknown',
                'source': source,
                'text': '',
                'date': None
            }
        
        try:
            # Get file extension
            file_ext = file_path.suffix.lower()
            
            # Check if format is supported
            if file_ext not in self.extractors:
                error_msg = f"Unsupported file format: {file_ext}"
                file_key = source or str(file_path)
                self.error_log[file_key] = error_msg
                
                # Determine filename based on source type
                if source and self._is_url(source):
                    filename = self._extract_filename_from_url(source)
                else:
                    filename = file_path.name
                
                return {
                    'filename': filename,
                    'source': source,
                    'text': '',
                    'date': None
                }
            
            # Extract content and metadata using appropriate extractor
            try:
                extractor = self.extractors[file_ext]
                content = extractor.extract_text(file_path)
                format_metadata = extractor.extract_metadata(file_path)
            except IngestionError as e:
                # Document-specific extraction failure (corrupted file, etc.)
                file_key = source or str(file_path)
                self.error_log[file_key] = f"Text extraction failed: {e}"
                
                if source and self._is_url(source):
                    filename = self._extract_filename_from_url(source)
                else:
                    filename = file_path.name
                
                return {
                    'filename': filename,
                    'source': source,
                    'text': '',
                    'date': None
                }
            except OSError as e:
                # File system access failure
                file_key = source or str(file_path)
                self.error_log[file_key] = f"File access failed: {e}"
                
                if source and self._is_url(source):
                    filename = self._extract_filename_from_url(source)
                else:
                    filename = file_path.name
                
                return {
                    'filename': filename,
                    'source': source,
                    'text': '',
                    'date': None
                }
            
            # Get file system metadata (modification date as fallback)
            try:
                stat = file_path.stat()
                modification_date = datetime.fromtimestamp(stat.st_mtime).isoformat()
            except OSError as e:
                # stat() failed - not critical, continue without file system metadata
                modification_date = None
            except ValueError as e:
                # Timestamp conversion failed (invalid timestamp, etc.)
                modification_date = None
            
            # Use document date if available, otherwise file modification date
            date_value = format_metadata.get('date') or modification_date
            
            # Determine filename based on source type
            if source and self._is_url(source):
                # Extract filename from URL
                filename = source.split('/')[-1]
            else:
                # Use actual file path name
                filename = file_path.name
            
            return {
                'filename': filename,
                'source': source,
                'text': content,
                'date': date_value
            }
            
        except Exception as e:
            # Catch-all for unexpected errors
            error_msg = f"Unexpected extraction error: {e}"
            file_key = source or str(file_path)
            self.error_log[file_key] = error_msg
            logger.error(f"Unexpected error extracting {file_key}: {e}")
            
            # Attempt to determine filename even in failure case
            try:
                if source and self._is_url(source):
                    filename = self._extract_filename_from_url(source)
                else:
                    filename = file_path.name
            except Exception:
                filename = 'unknown'
            
            return {
                'filename': filename,
                'source': source,
                'text': '',
                'date': None
            }

    def collect(self, source: str, progress_callback: Callable[[int, int], None] | None = None) -> CorpusType:
        """Collect and process documents from file path or URL.
        
        Automatically detects whether source is a URL or filesystem path, then
        discovers and extracts text from all accessible documents. This is the
        primary public method for document collection.
        
        Args:
            source: Either a filesystem path (file or directory) or web URL.
                Paths can be absolute or relative. URLs must include protocol.
            progress_callback: Optional callback function called after each document
                is processed. Receives (current_count, total_count) as arguments.
                For filesystem sources, total_count is known upfront. For web sources,
                total_count updates as documents are discovered.

        Returns:
            Dictionary mapping sequential integer IDs (starting from 0) to
            document dictionaries. Each document dict contains:
                - 'filename': Name of the file
                - 'source': Original source (URL or file path)
                - 'text': Extracted text content (empty string on failure)
            
            Returns empty dict if source is invalid or no documents are found.
            Check error_log attribute for details on any failures.
        
        Raises:
            ValueError: If source is empty or invalid type.
            
        Note:
            This method implements graceful degradation - individual document
            failures do not halt collection. Failed documents are included in
            results with empty text, and details are logged to error_log.
            
            For URL sources, documents are downloaded to temporary files which
            are automatically cleaned up after processing.
        
        Example:
            CodexCollector = CodexCollector()
            
            # Collect from filesystem
            corpus = CodexCollector.collect('/documents')
            
            # Collect from web
            corpus = CodexCollector.collect('https://example.com/docs')
            
            # Check for errors
            if CodexCollector.error_log:
                for source, error in CodexCollector.error_log.items():
                    print(f"Failed: {source} - {error}")
        """
        # Input validation
        if not isinstance(source, str):
            raise ValueError(
                f"source must be string, got {type(source).__name__}"
            )

        if not source or not source.strip():
            raise ValueError("source cannot be empty or whitespace-only")

        # Clear error log from previous collection
        self.error_log.clear()

        # Track collection start time for timeout
        start_time = time.time()

        results: dict[int, dict[str, str | None]] = {}
        index = 0
        
        try:
            if self._is_url(source):
                # Handle URL - discover and process documents
                discovered_urls = self._discover_files_from_url(source)
                
                if not discovered_urls:
                    logger.info(f"No documents discovered from URL: {source}")
                
                for url in discovered_urls:
                    # Check timeout
                    if self._check_collection_timeout(start_time):
                        logger.warning(f"Collection timeout reached ({self.max_collection_time}s)")
                        break

                    temp_file = None
                    try:
                        # Download file to temporary location with unique name
                        timestamp = datetime.now().timestamp()
                        temp_file = Path(f"temp_{timestamp}_{index}")
                        
                        # Extract filename from URL and preserve extension
                        url_filename = self._extract_filename_from_url(url)
                        if '.' in url_filename and url_filename != 'unknown':
                            extension = '.' + url_filename.split('.')[-1]
                            temp_file = temp_file.with_suffix(extension)
                        
                        self._download_file(url, temp_file)
                        
                        # Extract text from downloaded file, passing URL as source
                        result = self._extract_from_file(temp_file, url)
                        results[index] = result
                        index += 1

                        # Call progress callback if provided
                        if progress_callback:
                            progress_callback(index, len(discovered_urls))
                        
                    except IngestionError as e:
                        # Document extraction failed (already logged by _extract_from_file)
                        error_msg = f"Document processing failed: {e}"
                        self.error_log[url] = error_msg
                        results[index] = {
                            'filename': self._extract_filename_from_url(url),
                            'source': url,
                            'text': '',
                            'date': None
                        }
                        index += 1
                        
                    except OSError as e:
                        # File system errors (disk full, temp directory issues)
                        error_msg = f"Temporary file handling failed: {e}"
                        self.error_log[url] = error_msg
                        results[index] = {
                            'filename': self._extract_filename_from_url(url),
                            'source': url,
                            'text': '',
                            'date': None
                        }
                        index += 1
                        
                    except Exception as e:
                        # Unexpected errors
                        error_msg = f"Unexpected error processing URL: {e}"
                        self.error_log[url] = error_msg
                        logger.error(f"Unexpected error processing {url}: {e}")
                        results[index] = {
                            'filename': self._extract_filename_from_url(url),
                            'source': url,
                            'text': '',
                            'date': None
                        }
                        index += 1
                        
                    finally:
                        # Clean up temporary file
                        if temp_file and temp_file.exists():
                            try:
                                temp_file.unlink()
                            except OSError as e:
                                # Log cleanup failure but don't halt processing
                                logger.warning(f"Failed to delete temporary file {temp_file}: {e}")
            
            else:
                # Handle file path - discover and process files
                discovered_files = self._discover_files_from_path(source)
                
                if not discovered_files:
                    logger.info(f"No valid files discovered from path: {source}")
                
                total_files = len(discovered_files)
                for file_path_str in discovered_files:
                    # Check timeout
                    if self._check_collection_timeout(start_time):
                        logger.warning(f"Collection timeout reached ({self.max_collection_time}s)")
                        break

                    try:
                        file_path = Path(file_path_str)
                        # Pass file path string as source
                        result = self._extract_from_file(file_path, file_path_str)
                        results[index] = result
                        index += 1

                        # Call progress callback if provided
                        if progress_callback:
                            progress_callback(index, total_files)
                        
                    except (ValueError, OSError) as e:
                        # Path construction or file access failures
                        error_msg = f"File processing failed: {e}"
                        self.error_log[file_path_str] = error_msg
                        results[index] = {
                            'filename': Path(file_path_str).name,
                            'source': file_path_str,
                            'text': '',
                            'date': None
                        }
                        index += 1
                        
                    except Exception as e:
                        # Unexpected errors
                        error_msg = f"Unexpected error processing file: {e}"
                        self.error_log[file_path_str] = error_msg
                        logger.error(f"Unexpected error processing {file_path_str}: {e}")
                        results[index] = {
                            'filename': Path(file_path_str).name,
                            'source': file_path_str,
                            'text': '',
                            'date': None
                        }
                        index += 1
            
            self.corpus = results
            logger.info(f"Collection complete: {len(results)} documents, {len(self.error_log)} errors")
            return results
            
        except ValueError as e:
            # Input validation or _is_url failures
            error_msg = f"Invalid source: {e}"
            self.error_log[source] = error_msg
            self.corpus = results
            return results
            
        except Exception as e:
            # Unexpected top-level failure
            error_msg = f"Collection failed: {e}"
            self.error_log[source] = error_msg
            logger.error(f"Unexpected collection failure for {source}: {e}")
            self.corpus = results
            return results