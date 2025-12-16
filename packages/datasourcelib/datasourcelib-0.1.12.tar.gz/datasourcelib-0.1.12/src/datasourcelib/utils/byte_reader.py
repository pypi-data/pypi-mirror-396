from pathlib import Path
from typing import Optional, Union, List
import io
import pandas as pd

# --- Optional helpers ---
from charset_normalizer import from_bytes as cn_from_bytes

# DOCX
from docx import Document as DocxDocument

# PDF
import fitz  # pymupdf
import pdfplumber

# PPTX
from pptx import Presentation

# YAML / XML
import yaml
from lxml import etree
import json


class ByteReader:
    """
    Unified reader for common file types.
    - read_text(path): file path -> text
    - read_table(path): file path -> DataFrame
    - read_text_from_bytes(data, ext): bytes -> text
    - read_table_from_bytes(data, ext): bytes -> DataFrame
    """

    TEXT_EXTS = {".txt", ".log", ".md"}
    TABLE_EXTS = {".csv", ".tsv", ".xlsx", ".xls"}
    DOCX_EXTS = {".docx"}
    PDF_EXTS = {".pdf"}
    PPTX_EXTS = {".pptx"}
    JSON_EXTS = {".json"}
    YAML_EXTS = {".yaml", ".yml"}
    INI_EXTS = {".ini", ".cfg"}
    XML_EXTS = {".xml"}

    def __init__(self, default_encoding: str = "utf-8", errors: str = "replace"):
        self.default_encoding = default_encoding
        self.errors = errors

    # -----------------------
    # Public API (paths)
    # -----------------------
    def read_text(self, path: Union[str, Path]) -> str:
        path = Path(path)
        ext = path.suffix.lower()

        if ext in self.TEXT_EXTS:
            return path.read_text(encoding=self.default_encoding, errors=self.errors)

        if ext in self.PDF_EXTS:
            return self._read_pdf_text_path(path)

        if ext in self.DOCX_EXTS:
            return self._read_docx_text_fp(open(path, "rb"))

        if ext in self.PPTX_EXTS:
            return self._read_pptx_text_fp(open(path, "rb"))

        if ext in self.JSON_EXTS:
            with path.open("r", encoding=self.default_encoding, errors=self.errors) as f:
                obj = json.load(f)
            return json.dumps(obj, indent=2, ensure_ascii=False)

        if ext in self.YAML_EXTS:
            with path.open("r", encoding=self.default_encoding, errors=self.errors) as f:
                obj = yaml.safe_load(f)
            return yaml.safe_dump(obj, sort_keys=False, allow_unicode=True)

        if ext in self.INI_EXTS:
            import configparser
            parser = configparser.ConfigParser()
            with path.open("r", encoding=self.default_encoding, errors=self.errors) as f:
                parser.read_file(f)
            output = io.StringIO()
            parser.write(output)
            return output.getvalue()

        if ext in self.XML_EXTS:
            tree = etree.parse(str(path))
            return etree.tostring(tree, pretty_print=True, encoding="unicode")

        if ext in self.TABLE_EXTS:
            df = self.read_table(path)
            return df.to_csv(index=False)

        raise ValueError(f"Unsupported file extension for text extraction: {ext}")

    def read_table(self, path: Union[str, Path]) -> pd.DataFrame:
        path = Path(path)
        ext = path.suffix.lower()

        if ext == ".csv":
            return pd.read_csv(path)
        if ext == ".tsv":
            return pd.read_csv(path, sep="\t")
        if ext == ".xlsx":
            return pd.read_excel(path, engine="openpyxl")
        if ext == ".xls":
            return pd.read_excel(path, engine="xlrd")

        # Fallback: attempt CSV read if unknown
        try:
            return pd.read_csv(path)
        except Exception as e:
            raise ValueError(f"Unsupported file extension for tables: {ext}") from e

    # -----------------------
    # Public API (bytes)
    # -----------------------
    def read_text_from_bytes(self, data: bytes, ext: str) -> str:
        """
        Extract text from in-memory bytes.
        ext: file extension (e.g., '.pdf', '.docx', '.txt', '.pptx', '.json', '.yaml', '.xml', '.csv', '.xlsx')
        """
        ext = self._normalize_ext(ext)

        if ext in self.TEXT_EXTS:
            # Robust encoding detection
            res = cn_from_bytes(data).best()
            return str(res) if res else data.decode(self.default_encoding, errors=self.errors)

        if ext in self.PDF_EXTS:
            return self._read_pdf_text_bytes(data)

        if ext in self.DOCX_EXTS:
            return self._read_docx_text_fp(io.BytesIO(data))

        if ext in self.PPTX_EXTS:
            return self._read_pptx_text_fp(io.BytesIO(data))

        if ext in self.JSON_EXTS:
            obj = json.loads(data.decode(self.default_encoding, errors=self.errors))
            return json.dumps(obj, indent=2, ensure_ascii=False)

        if ext in self.YAML_EXTS:
            obj = yaml.safe_load(data.decode(self.default_encoding, errors=self.errors))
            return yaml.safe_dump(obj, sort_keys=False, allow_unicode=True)

        if ext in self.INI_EXTS:
            import configparser
            parser = configparser.ConfigParser()
            parser.read_string(data.decode(self.default_encoding, errors=self.errors))
            output = io.StringIO()
            parser.write(output)
            return output.getvalue()

        if ext in self.XML_EXTS:
            tree = etree.parse(io.BytesIO(data))
            return etree.tostring(tree, pretty_print=True, encoding="unicode")

        if ext in self.TABLE_EXTS:
            df = self.read_table_from_bytes(data, ext)
            return df.to_csv(index=False)

        raise ValueError(f"Unsupported extension for text extraction from bytes: {ext}")

    def read_table_from_bytes(self, data: bytes, ext: str) -> pd.DataFrame:
        """
        Load tabular data from in-memory bytes into a DataFrame.
        """
        ext = self._normalize_ext(ext)

        if ext == ".csv":
            return pd.read_csv(io.BytesIO(data))
        if ext == ".tsv":
            return pd.read_csv(io.BytesIO(data), sep="\t")
        if ext == ".xlsx":
            return pd.read_excel(io.BytesIO(data), engine="openpyxl")
        if ext == ".xls":
            return pd.read_excel(io.BytesIO(data), engine="xlrd")

        # Opportunistic fallback: try CSV
        try:
            return pd.read_csv(io.BytesIO(data))
        except Exception as e:
            raise ValueError(f"Unsupported extension for table reading from bytes: {ext}") from e

    # -----------------------
    # Internal helpers
    # -----------------------
    def _normalize_ext(self, ext: str) -> str:
        ext = (ext or "").strip().lower()
        if not ext.startswith("."):
            ext = "." + ext
        return ext

    def _read_pdf_text_path(self, path: Path) -> str:
        # Prefer PyMuPDF
        try:
            parts: List[str] = []
            with fitz.open(str(path)) as doc:
                if doc.is_encrypted and not doc.authenticate(""):
                    raise RuntimeError("Encrypted PDF requires a password.")
                for page in doc:
                    parts.append(page.get_text("text"))
            text = "\n\n".join(parts).strip()
            if text:
                return text
        except Exception:
            pass

        # Fallback: pdfplumber
        with pdfplumber.open(str(path)) as pdf:
            return "\n\n".join([(p.extract_text() or "") for p in pdf.pages]).strip()

    def _read_pdf_text_bytes(self, data: bytes) -> str:
        # PyMuPDF can open from bytes
        try:
            doc = fitz.open(stream=data, filetype="pdf")
            parts: List[str] = []
            if doc.is_encrypted and not doc.authenticate(""):
                raise RuntimeError("Encrypted PDF requires a password.")
            for page in doc:
                parts.append(page.get_text("text"))
            doc.close()
            text = "\n\n".join(parts).strip()
            if text:
                return text
        except Exception:
            pass

        # Fallback to pdfplumber from BytesIO
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            return "\n\n".join([(p.extract_text() or "") for p in pdf.pages]).strip()

    def _read_docx_text_fp(self, fp) -> str:
        doc = DocxDocument(fp)
        chunks = []
        for p in doc.paragraphs:
            if p.text:
                chunks.append(p.text)
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                if any(cells):
                    chunks.append("\t".join(cells))
        return "\n".join(chunks).strip()

    def _read_pptx_text_fp(self, fp) -> str:
        prs = Presentation(fp)
        chunks = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    text = shape.text or ""
                    if text:
                        chunks.append(text)
        return "\n".join(chunks).strip()
