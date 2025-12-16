from pathlib import Path
from typing import Optional, Union, List
import io
import pandas as pd

# --- Optional helpers ---
from charset_normalizer import from_path as cn_from_path

# DOCX
from docx import Document as DocxDocument

# PDF
import fitz  # pymupdf
import pdfplumber

# PPTX
from pptx import Presentation

# YAML / XML
import yaml
from lxml import etree


class FileReader:
    """
    A unified reader for common file types.
    - read_text(path): extract text from txt, pdf, docx, pptx, json, yaml, ini, xml
    - read_table(path): load tabular data from csv, tsv, xlsx/xls
    """

    TEXT_EXTS = {".txt", ".log", ".md"}
    TABLE_EXTS = {".csv", ".tsv", ".xlsx", ".xls"}
    DOCX_EXTS = {".docx"}
    PDF_EXTS = {".pdf"}
    PPTX_EXTS = {".pptx"}
    JSON_EXTS = {".json"}
    YAML_EXTS = {".yaml", ".yml"}
    INI_EXTS = {".ini", ".cfg"}
    XML_EXTS = {".xml"}

    def __init__(self, default_encoding: str = "utf-8", errors: str = "replace"):
        self.default_encoding = default_encoding
        self.errors = errors

    # -----------------------
    # Public API
    # -----------------------
    def read_text(self, path: Union[str, Path]) -> str:
        """
        Extract best-effort text from a given file based on extension.
        """
        path = Path(path)
        ext = path.suffix.lower()

        if ext in self.TEXT_EXTS:
            return self._read_plain_text(path)

        if ext in self.PDF_EXTS:
            return self._read_pdf_text(path)

        if ext in self.DOCX_EXTS:
            return self._read_docx_text(path)

        if ext in self.PPTX_EXTS:
            return self._read_pptx_text(path)

        if ext in self.JSON_EXTS:
            return self._read_json_text(path)

        if ext in self.YAML_EXTS:
            return self._read_yaml_text(path)

        if ext in self.INI_EXTS:
            return self._read_ini_text(path)

        if ext in self.XML_EXTS:
            return self._read_xml_text(path)

        if ext in self.TABLE_EXTS:
            # For tabular files, provide a quick text representation
            df = self.read_table(path)
            return df.to_csv(index=False)

        raise ValueError(f"Unsupported file extension for text extraction: {ext}")

    def read_table(self, path: Union[str, Path]) -> pd.DataFrame:
        """
        Load tabular data from CSV/TSV/Excel, returning a DataFrame.
        """
        path = Path(path)
        ext = path.suffix.lower()

        if ext == ".csv":
            return pd.read_csv(path)
        if ext == ".tsv":
            return pd.read_csv(path, sep="\t")
        if ext == ".xlsx":
            return pd.read_excel(path, engine="openpyxl")
        if ext == ".xls":
            return pd.read_excel(path, engine="xlrd")

        # Fallback: attempt CSV read if unknown
        try:
            return pd.read_csv(path)
        except Exception as e:
            raise ValueError(f"Unsupported file extension for tables: {ext}") from e

    # -----------------------
    # Text readers
    # -----------------------
    def _read_plain_text(self, path: Path) -> str:
        # Detect encoding for robustness
        res = cn_from_path(str(path)).best()
        if res:
            return str(res)
        # Fallback to configured defaults
        return path.read_text(encoding=self.default_encoding, errors=self.errors)

    def _read_pdf_text(self, path: Path) -> str:
        # Try PyMuPDF (fast, layout-aware)
        try:
            text_parts: List[str] = []
            with fitz.open(str(path)) as doc:
                if doc.can_save and doc.is_encrypted:
                    # If encrypted and requires a password, this will fail to extract text.
                    if not doc.authenticate(""):
                        raise RuntimeError("Encrypted PDF requires a password.")
                for page in doc:
                    text_parts.append(page.get_text("text"))
            text = "\n".join(text_parts).strip()
            if text:
                return text
        except Exception:
            pass

        # Fallback to pdfplumber (good for tables/structured text)
        try:
            text_parts = []
            with pdfplumber.open(str(path)) as pdf:
                for page in pdf.pages:
                    t = page.extract_text() or ""
                    text_parts.append(t)
            return "\n".join(text_parts).strip()
        except Exception as e:
            raise RuntimeError(f"Failed to read PDF: {e}") from e

    def _read_docx_text(self, path: Path) -> str:
        doc = DocxDocument(str(path))
        chunks = []
        # Paragraphs
        for p in doc.paragraphs:
            if p.text:
                chunks.append(p.text)
        # Tables (optional: include)
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                if any(cells):
                    chunks.append("\t".join(cells))
        return "\n".join(chunks).strip()

    def _read_pptx_text(self, path: Path) -> str:
        prs = Presentation(str(path))
        chunks = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.has_text_frame:
                    text = shape.text if hasattr(shape, "text") else ""
                    if text:
                        chunks.append(text)
        return "\n".join(chunks).strip()

    def _read_json_text(self, path: Path) -> str:
        import json
        with path.open("r", encoding=self.default_encoding, errors=self.errors) as f:
            obj = json.load(f)
        # Pretty-print
        return json.dumps(obj, indent=2, ensure_ascii=False)

    def _read_yaml_text(self, path: Path) -> str:
        with path.open("r", encoding=self.default_encoding, errors=self.errors) as f:
            obj = yaml.safe_load(f)
        return yaml.safe_dump(obj, sort_keys=False, allow_unicode=True)

    def _read_ini_text(self, path: Path) -> str:
        import configparser
        parser = configparser.ConfigParser()
        with path.open("r", encoding=self.default_encoding, errors=self.errors) as f:
            # INI files might have duplicate keys; defaults handle many cases
            parser.read_file(f)
        output = io.StringIO()
        parser.write(output)
        return output.getvalue()

    def _read_xml_text(self, path: Path) -> str:
        # Pretty-print XML
        tree = etree.parse(str(path))
        return etree.tostring(tree, pretty_print=True, encoding="unicode")


# -----------------------
# Example usage
# -----------------------
#if __name__ == "__main__":
#    reader = FileReader()

    # 1) Extract text
    # print(reader.read_text("document.pdf"))
    # print(reader.read_text("report.docx"))
    # print(reader.read_text("slides.pptx"))
    # print(reader.read_text("notes.txt"))
    # print(reader.read_text("config.yaml"))
    # print(reader.read_text("data.xml"))

    # 2) Load tabular data
    # df = reader.read_table("data.xlsx")
    # print(df.head())
